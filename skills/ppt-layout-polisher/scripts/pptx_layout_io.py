from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
EMU_PER_INCH = 914400.0


def emu_to_inches(value: int | float) -> float:
    return round(float(value) / EMU_PER_INCH, 4)


def pptx_to_document(path: str | Path, margin: float = 0.6) -> dict:
    prs = Presentation(str(path))
    page = {
        "width": round(prs.slide_width / EMU_PER_INCH, 3),
        "height": round(prs.slide_height / EMU_PER_INCH, 3),
        "unit": "in",
        "margin": margin,
    }
    slides = []
    for slide_index, slide in enumerate(prs.slides):
        slides.append(
            {
                "id": f"slide-{slide_index + 1}",
                "type": infer_slide_type(slide),
                "title": slide.shapes.title.text if slide.shapes.title and slide.shapes.title.has_text_frame else "",
                "elements": [shape_to_element(shape, shape_index) for shape_index, shape in enumerate(slide.shapes)],
            }
        )
    return {"page": page, "slides": slides}


def export_layout_json(pptx_path: str | Path, output_path: str | Path, margin: float = 0.6) -> None:
    Path(output_path).write_text(
        json.dumps(pptx_to_document(pptx_path, margin=margin), ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


def infer_slide_type(slide) -> str:
    title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.has_text_frame else ""
    text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
    picture_shapes = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    table_shapes = [shape for shape in slide.shapes if getattr(shape, "has_table", False)]

    if slide.shapes.title and len(slide.shapes) <= 3:
        if len(text_shapes) <= 2 and not picture_shapes:
            return "section" if len(text_shapes) <= 1 else "cover"
    if table_shapes:
        return "table"
    if len(picture_shapes) == 1 and len(text_shapes) >= 2:
        return "image-text"
    if len(text_shapes) == 2:
        return "comparison"
    if len(text_shapes) == 3:
        return "cards"
    if title:
        return "single-column"
    return "content"


def shape_to_element(shape, shape_index: int) -> dict:
    text = extract_text(shape)
    font_family, font_size, line_height = extract_font(shape)
    kind = infer_shape_kind(shape)
    role = infer_text_role(shape, kind, text, font_size)

    return {
        "id": f"shape-{shape_index}",
        "kind": kind,
        "text_role": role,
        "x": emu_to_inches(shape.left),
        "y": emu_to_inches(shape.top),
        "w": emu_to_inches(shape.width),
        "h": emu_to_inches(shape.height),
        "text": text,
        "font_family": font_family,
        "font_size": font_size,
        "line_height": line_height,
        "padding": 0.08 if text else 0.0,
    }


def extract_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    lines: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        line = "".join(run.text for run in paragraph.runs).strip()
        if not line and paragraph.text:
            line = paragraph.text.strip()
        lines.append(line)
    return "\n".join(line for line in lines if line).strip()


def extract_font(shape) -> tuple[str, float, float]:
    if not getattr(shape, "has_text_frame", False):
        return "", 0.0, 0.0
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = run.text.strip()
            if not text:
                continue
            font = run.font
            size = round(font.size.pt, 2) if font.size else 0.0
            return font.name or "", size, 1.2
    return "", 0.0, 1.2


def infer_shape_kind(shape) -> str:
    if getattr(shape, "has_table", False):
        return "table"
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if getattr(shape, "has_chart", False):
        return "chart"
    if getattr(shape, "has_text_frame", False):
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            if placeholder_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}:
                return "title"
            if placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                return "subtitle"
        return "text"
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return "shape"
    return "group" if shape.shape_type == MSO_SHAPE_TYPE.GROUP else "shape"


def infer_text_role(shape, kind: str, text: str, font_size: float) -> str:
    if not text:
        return "none"
    if shape.is_placeholder:
        placeholder_type = shape.placeholder_format.type
        if placeholder_type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}:
            return "title"
        if placeholder_type == PP_PLACEHOLDER.SUBTITLE:
            return "subtitle"
        if placeholder_type == PP_PLACEHOLDER.BODY:
            return "body"
    top = emu_to_inches(shape.top)
    width = emu_to_inches(shape.width)
    if top <= 1.2 and width >= 5.5 and len(text.replace("\n", " ")) <= 90:
        return "title"
    if kind == "title" or font_size >= 24:
        return "title"
    if font_size <= 12 and emu_to_inches(shape.top) >= 5.5:
        return "note"
    return "body"
