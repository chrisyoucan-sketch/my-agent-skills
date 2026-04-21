from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

from detect_text_overlap import detect_text_overlap
from inspect_ppt_layout import inspect_doc
from layout_common import SAFE_CHINESE_FONTS, SAFE_LATIN_FONTS, estimate_text_capacity
from pptx_layout_io import EMU_PER_INCH, emu_to_inches, pptx_to_document, shape_to_element


SAFE_FONTS = SAFE_CHINESE_FONTS | SAFE_LATIN_FONTS


def contains_cjk(text: str) -> bool:
    return any("\u4e00" <= ch <= "\u9fff" for ch in text)


def target_font(text: str) -> str:
    return "Microsoft YaHei" if contains_cjk(text) else "Aptos"


def role_defaults(role: str) -> tuple[float, float]:
    if role == "title":
        return 28.0, 20.0
    if role == "subtitle":
        return 18.0, 14.0
    if role == "note":
        return 11.0, 10.0
    return 16.0, 12.0


def slide_elements(slide) -> list[dict]:
    return [shape_to_element(shape, index) for index, shape in enumerate(slide.shapes)]


def clamp_shape(shape, slide_width: int, slide_height: int, margin: int) -> bool:
    changed = False
    max_width = max(slide_width - 2 * margin, int(1.0 * EMU_PER_INCH))
    max_height = max(slide_height - 2 * margin, int(0.5 * EMU_PER_INCH))
    if shape.width > max_width:
        shape.width = max_width
        changed = True
    if shape.height > max_height:
        shape.height = max_height
        changed = True

    max_left = max(margin, slide_width - margin - shape.width)
    max_top = max(margin, slide_height - margin - shape.height)
    new_left = min(max(shape.left, margin), max_left)
    new_top = min(max(shape.top, margin), max_top)
    if new_left != shape.left:
        shape.left = new_left
        changed = True
    if new_top != shape.top:
        shape.top = new_top
        changed = True
    return changed


def snap_attr(shapes: list, attr: str, tolerance_in: float) -> int:
    tol = int(tolerance_in * EMU_PER_INCH)
    ordered = sorted(shapes, key=lambda shape: getattr(shape, attr))
    clusters: list[list] = []
    for shape in ordered:
        if not clusters or getattr(shape, attr) - getattr(clusters[-1][-1], attr) > tol:
            clusters.append([shape])
        else:
            clusters[-1].append(shape)

    changed = 0
    for cluster in clusters:
        if len(cluster) < 2:
            continue
        anchor = int(round(sum(getattr(shape, attr) for shape in cluster) / len(cluster)))
        for shape in cluster:
            if getattr(shape, attr) != anchor:
                setattr(shape, attr, anchor)
                changed += 1
    return changed


def normalize_fonts(shape, role: str) -> int:
    if not getattr(shape, "has_text_frame", False):
        return 0

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    default_pt, min_pt = role_defaults(role)
    changed = 0

    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            text = run.text or ""
            if not text.strip():
                continue
            font = run.font
            desired = target_font(text)
            if not font.name or font.name not in SAFE_FONTS:
                font.name = desired
                changed += 1
            current = font.size.pt if font.size else 0.0
            if current <= 0:
                font.size = Pt(default_pt)
                changed += 1
            elif current < min_pt:
                font.size = Pt(min_pt)
                changed += 1
    return changed


def estimate_shape_overflow(shape, index: int) -> bool:
    element = shape_to_element(shape, index)
    if not element["text"]:
        return False
    return estimate_text_capacity(element)["overflow"] > 0


def reduce_font_size(shape, role: str) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    _, min_pt = role_defaults(role)
    changed = False
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text.strip() or not run.font.size:
                continue
            if run.font.size.pt > min_pt:
                run.font.size = Pt(max(min_pt, run.font.size.pt - 1))
                changed = True
    return changed


def grow_height(shape, slide_height: int, margin: int, max_growth_in: float = 0.35) -> bool:
    available = slide_height - margin - shape.top - shape.height
    growth = min(available, int(max_growth_in * EMU_PER_INCH))
    if growth <= 0:
        return False
    shape.height += growth
    return True


def relayout_as_two_columns(left, right, slide_width: int, margin: int, gap: int) -> bool:
    usable = slide_width - 2 * margin - gap
    if usable <= int(3.0 * EMU_PER_INCH):
        return False
    col_width = int(usable / 2)
    top = min(left.top, right.top)
    changed = False

    if left.left != margin:
        left.left = margin
        changed = True
    if right.left != margin + col_width + gap:
        right.left = margin + col_width + gap
        changed = True
    if left.width != col_width:
        left.width = col_width
        changed = True
    if right.width != col_width:
        right.width = col_width
        changed = True
    if left.top != top:
        left.top = top
        changed = True
    if right.top != top:
        right.top = top
        changed = True
    return changed


def fit_text_shapes(slide, slide_height: int, margin: int) -> int:
    changed = 0
    for index, shape in enumerate(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        role = shape_to_element(shape, index)["text_role"]
        if not estimate_shape_overflow(shape, index):
            continue
        if grow_height(shape, slide_height, margin):
            changed += 1
        for _ in range(4):
            if not estimate_shape_overflow(shape, index):
                break
            if not reduce_font_size(shape, role):
                break
            changed += 1
    return changed


def resolve_text_overlaps(slide, slide_width: int, slide_height: int, margin: int) -> int:
    changed = 0
    gap = int(0.06 * EMU_PER_INCH)
    for _ in range(12):
        overlaps = detect_text_overlap({"elements": slide_elements(slide)}, None)
        if not overlaps:
            break
        any_change = False
        for finding in overlaps:
            left = slide.shapes[int(finding["left"].split("-")[1])]
            right = slide.shapes[int(finding["right"].split("-")[1])]
            pair_left, pair_right = sorted([left, right], key=lambda shape: shape.left)
            horizontal_overlap = pair_left.left + pair_left.width + gap > pair_right.left
            similar_band = abs(pair_left.top - pair_right.top) <= int(0.5 * EMU_PER_INCH)
            if horizontal_overlap and similar_band:
                if relayout_as_two_columns(pair_left, pair_right, slide_width, margin, gap):
                    any_change = True
                    changed += 1
                    continue

            upper, lower = sorted([left, right], key=lambda shape: (shape.top, shape.left))
            new_top = upper.top + upper.height + gap
            if new_top + lower.height <= slide_height - margin:
                lower.top = new_top
                any_change = True
                changed += 1
                continue
            new_left = upper.left + upper.width + gap
            if new_left + lower.width <= slide_width - margin:
                lower.left = new_left
                any_change = True
                changed += 1
                continue
            if grow_height(upper, slide_height, margin, max_growth_in=0.15):
                any_change = True
                changed += 1
        if not any_change:
            break
    return changed


def fix_slide(slide, slide_width: int, slide_height: int, margin_in: float) -> int:
    margin = int(margin_in * EMU_PER_INCH)
    changed = 0
    text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]

    for index, shape in enumerate(slide.shapes):
        role = shape_to_element(shape, index)["text_role"]
        if getattr(shape, "has_text_frame", False):
            changed += normalize_fonts(shape, role)
        changed += int(clamp_shape(shape, slide_width, slide_height, margin))

    changed += snap_attr(text_shapes, "left", tolerance_in=0.08)
    changed += snap_attr(text_shapes, "top", tolerance_in=0.08)

    for shape in slide.shapes:
        changed += int(clamp_shape(shape, slide_width, slide_height, margin))

    changed += resolve_text_overlaps(slide, slide_width, slide_height, margin)
    changed += fit_text_shapes(slide, slide_height, margin)

    for shape in slide.shapes:
        changed += int(clamp_shape(shape, slide_width, slide_height, margin))
    return changed


def auto_fix_presentation(pptx_path: str, output_path: str, margin: float) -> dict:
    before = inspect_doc(pptx_to_document(pptx_path, margin=margin))
    prs = Presentation(pptx_path)
    changes = 0
    for slide in prs.slides:
        changes += fix_slide(slide, prs.slide_width, prs.slide_height, margin)
    prs.save(output_path)
    after = inspect_doc(pptx_to_document(output_path, margin=margin))
    return {"changes": changes, "before": before, "after": after}


def main() -> None:
    parser = argparse.ArgumentParser(description="Auto-fix common PPT layout issues directly in a PPTX file.")
    parser.add_argument("pptx_path", help="Path to the input .pptx file.")
    parser.add_argument("-o", "--output", required=True, help="Path to the output .pptx file.")
    parser.add_argument("--margin", type=float, default=0.6, help="Safety margin in inches. Default: 0.6")
    parser.add_argument("--report", help="Optional JSON report path.")
    args = parser.parse_args()

    report = auto_fix_presentation(args.pptx_path, args.output, args.margin)
    payload = json.dumps(report, ensure_ascii=False, indent=2)
    if args.report:
        Path(args.report).write_text(payload, encoding="utf-8")
    else:
        print(payload)


if __name__ == "__main__":
    main()
